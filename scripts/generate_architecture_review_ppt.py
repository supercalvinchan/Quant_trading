#!/usr/bin/env python3
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "USalpha_AI_Agent_Architecture_Review_v4.pptx"


BG = RGBColor(14, 18, 28)
PANEL = RGBColor(25, 31, 45)
PANEL_2 = RGBColor(33, 40, 58)
ACCENT = RGBColor(57, 181, 255)
ACCENT_2 = RGBColor(129, 211, 248)
GREEN = RGBColor(87, 214, 157)
YELLOW = RGBColor(255, 196, 87)
RED = RGBColor(255, 120, 120)
WHITE = RGBColor(243, 247, 250)
MUTED = RGBColor(176, 188, 204)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(11.9), Inches(0.7))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.58), Inches(0.95), Inches(11.5), Inches(0.4))
        stf = sbox.text_frame
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        sr.font.size = Pt(10.5)
        sr.font.color.rgb = MUTED


def add_footer(slide, page: str):
    box = slide.shapes.add_textbox(Inches(10.9), Inches(7.0), Inches(1.6), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = page
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def add_panel(slide, left, top, width, height, title=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = PANEL_2
    shape.line.width = Pt(1.0)
    if title:
        tbox = slide.shapes.add_textbox(left + Inches(0.18), top + Inches(0.08), width - Inches(0.3), Inches(0.3))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = ACCENT_2
    return shape


def add_bullets(slide, left, top, width, height, items, font_size=14, color=WHITE, level_step=0.22):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(2)
    tf.margin_bottom = Pt(2)
    first = True
    for item in items:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        p.text = text
        p.font.size = Pt(font_size - level * 1.2)
        p.font.color.rgb = color
        p.space_after = Pt(4)
        if level == 0:
            p.bullet = True
        p.left_margin = Inches(level * level_step)
    return box


def add_kpi(slide, left, top, width, height, title, value, color):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PANEL
    shape.line.color.rgb = PANEL_2
    tb = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), height - Inches(0.16))
    tf = tb.text_frame
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.LEFT
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(10.5)
    r1.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = value
    r2.font.size = Pt(22)
    r2.font.bold = True
    r2.font.color.rgb = color


def add_table(slide, left, top, width, height, headers, rows, col_widths=None, font_size=10.5):
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, head in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL_2
        cell.text = head
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(font_size)
                r.font.color.rgb = WHITE
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(font_size)
                    r.font.color.rgb = WHITE if j == 0 else MUTED
    return table


def add_flow_steps(
    slide,
    steps,
    *,
    left,
    top,
    width,
    height,
    gap=0.12,
    title_size=11,
    body_size=8.8,
    colors=None,
):
    n = len(steps)
    box_w = (width - Inches(gap * (n - 1))) / n
    for idx, step in enumerate(steps):
        title, body = step[:2]
        color = colors[idx] if colors and idx < len(colors) else PANEL
        x = left + idx * (box_w + Inches(gap))
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, top, box_w, height)
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.color.rgb = PANEL_2
        tb = slide.shapes.add_textbox(x + Inches(0.08), top + Inches(0.10), box_w - Inches(0.16), height - Inches(0.18))
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = title
        r1.font.size = Pt(title_size)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = WHITE
        if idx < n - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                x + box_w + Inches(0.015),
                top + height / 2 - Inches(0.13),
                Inches(0.09),
                Inches(0.26),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = ACCENT_2
            arrow.line.fill.background()


def overwrite_footer(slide, page: str):
    cover = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(10.82), Inches(6.93), Inches(1.8), Inches(0.32))
    cover.fill.solid()
    cover.fill.fore_color.rgb = BG
    cover.line.fill.background()
    add_footer(slide, page)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = BG
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(0.9), Inches(0.18), Inches(2.4))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()
    title = slide.shapes.add_textbox(Inches(1.05), Inches(0.95), Inches(10.8), Inches(1.8))
    tf = title.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "USalpha 架构复盘与\nAI Agent 量化平台路线图"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    sub = slide.shapes.add_textbox(Inches(1.08), Inches(3.0), Inches(10.5), Inches(0.8))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "现状评估 | 参考平台对比 | 功能缺口 | 模块化设计建议"
    sr.font.size = Pt(15)
    sr.font.color.rgb = ACCENT_2
    notes = slide.shapes.add_textbox(Inches(1.08), Inches(3.7), Inches(10.5), Inches(1.0))
    ntf = notes.text_frame
    np = ntf.paragraphs[0]
    nr = np.add_run()
    nr.text = "对象：USalpha 当前版本（v1）\n参考仓库：QuantaAlpha / alpha_mining / TradingAgents / akquant / a-stock-data / qlib / Kronos / vnpy"
    nr.font.size = Pt(12)
    nr.font.color.rgb = MUTED
    add_footer(slide, "01")


def slide_exec_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "执行摘要", "对当前 USalpha 的一句话判断与优先级建议")
    add_kpi(slide, Inches(0.6), Inches(1.35), Inches(2.0), Inches(1.0), "研究原型成熟度", "B", ACCENT)
    add_kpi(slide, Inches(2.8), Inches(1.35), Inches(2.0), Inches(1.0), "模块化程度", "B-", YELLOW)
    add_kpi(slide, Inches(5.0), Inches(1.35), Inches(2.0), Inches(1.0), "Agent 就绪度", "C", RED)
    add_kpi(slide, Inches(7.2), Inches(1.35), Inches(2.2), Inches(1.0), "实盘可用度", "C-", RED)
    add_panel(slide, Inches(0.6), Inches(2.6), Inches(5.9), Inches(3.8), "关键判断")
    add_bullets(
        slide,
        Inches(0.8),
        Inches(2.95),
        Inches(5.5),
        Inches(3.2),
        [
            "USalpha 已经具备“数据→因子→模型→回测→可视化”的完整研究闭环。",
            "它适合作为研究原型与产品雏形，但还不是成熟的平台型框架。",
            "当前最强的是：自包含、迭代快、功能链路短、容易改。",
            "当前最弱的是：数据层治理、回测引擎抽象、实验管理、Agent 编排与实盘桥接。 ",
        ],
    )
    add_panel(slide, Inches(6.7), Inches(2.6), Inches(5.95), Inches(3.8), "优先级建议")
    add_bullets(
        slide,
        Inches(6.9),
        Inches(2.95),
        Inches(5.5),
        Inches(3.2),
        [
            "短期：把 USalpha 从“能跑”变成“可持续研发”，先做数据层、缓存层、实验层治理。",
            "中期：引入统一策略/回测抽象、模型注册、评估注册、Agent 调度接口。",
            "长期：拆出 Agent Platform、Research Platform、Execution Platform 三层架构。",
        ],
    )
    add_footer(slide, "02")


def slide_usalpha_done(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "USalpha 已完成工作", "当前版本已经实现的功能面")
    add_table(
        slide,
        Inches(0.6),
        Inches(1.35),
        Inches(12.1),
        Inches(5.4),
        ["模块", "已完成内容", "评价"],
        [
            ["数据", "行情拉取、NASDAQ 列表缓存、A 股缓存、fallback 机制", "可用，但治理较浅"],
            ["因子", "内置 526 因子解析与计算，表达式 evaluator", "自包含是优点"],
            ["模型", "轻量截面收益预测基线", "适合快速实验"],
            ["回测", "多空回测 + A 股页面化回测", "功能闭环已形成"],
            ["LLM", "候选因子生成、打分、入库、结果留档", "有明显产品雏形"],
            ["可视化", "股票查看、训练挖掘、回测页面", "交互面已出现"],
            ["工程", "相对路径、统一入口、portable check、长任务稳定性修复", "工程意识在增强"],
        ],
        col_widths=[Inches(1.4), Inches(7.7), Inches(2.8)],
        font_size=11,
    )
    add_footer(slide, "03")


def slide_usalpha_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "USalpha 当前完整流程", "它目前完成的是“研究到候选交易”的流程，不是完整交易平台流程")
    steps = [
        ("1 股票池", "NASDAQ 列表\nA股全市场/自选"),
        ("2 行情准备", "yfinance / akshare\n本地缓存"),
        ("3 特征构建", "526 因子\n表达式计算"),
        ("4 预测打分", "截面收益预测\npred score"),
        ("5 组合回测", "Top/Bottom 分组\nlong-short"),
        ("6 候选输出", "明日候选股\n图表和持仓明细"),
        ("7 人工执行", "页面查看后\n手动交易"),
    ]
    add_flow_steps(
        slide,
        steps,
        left=Inches(0.45),
        top=Inches(1.9),
        width=Inches(12.35),
        height=Inches(1.55),
        gap=0.08,
        title_size=10.5,
        body_size=8.5,
        colors=[PANEL, PANEL, PANEL, PANEL, PANEL, RGBColor(45, 82, 120), RGBColor(90, 66, 38)],
    )
    add_panel(slide, Inches(0.6), Inches(4.0), Inches(5.9), Inches(2.2), "我们已经有")
    add_bullets(slide, Inches(0.82), Inches(4.35), Inches(5.45), Inches(1.55), [
        "股票池基础入口、行情缓存、因子计算、模型训练、回测结果、候选列表、页面展示。",
        "LLM 因子演进是研究增强模块，不是交易主流程的核心执行层。"
    ], font_size=11.5)
    add_panel(slide, Inches(6.72), Inches(4.0), Inches(5.95), Inches(2.2), "我们还没有")
    add_bullets(slide, Inches(6.94), Inches(4.35), Inches(5.45), Inches(1.55), [
        "组合构建规则、风控审批、订单管理、撮合/网关、成交回报、持仓同步、监控告警、复盘记忆。",
        "也就是说，USalpha 目前停在“研究结果 -> 人工看盘”的位置。"
    ], font_size=11.5)
    add_footer(slide, "04")


def slide_modularity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "模块化评价", "核心模块已经分开，但平台级抽象还不够")
    add_table(
        slide,
        Inches(0.55),
        Inches(1.35),
        Inches(12.2),
        Inches(5.6),
        ["维度", "现状", "评价", "建议"],
        [
            ["数据层", "US/CN 逻辑并存，采集/缓存/fallback 混在同层", "中等", "拆成 source / cache / catalog / quality"],
            ["因子层", "表达式与计算主干清楚", "较好", "补注册表、元数据、依赖关系"],
            ["模型层", "基线模型独立，但模型族很薄", "中等", "做 model registry 与训练接口"],
            ["回测层", "主流程回测与页面回测口径不统一", "偏弱", "统一 backtest engine / execution semantics"],
            ["页面层", "能力丰富，但单文件偏重", "偏弱", "拆 page / service / state / adapter"],
            ["实验层", "结果有 artifacts，但缺 experiment registry", "偏弱", "补 run metadata、comparison、repro"],
            ["Agent 层", "只有 LLM factor round 雏形", "较弱", "拆 planner / proposer / evaluator / reviewer / operator"],
        ],
        col_widths=[Inches(1.35), Inches(4.0), Inches(1.25), Inches(5.2)],
        font_size=10.5,
    )
    add_footer(slide, "05")


def slide_market_map(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "参考平台能力地图", "这些仓库分别代表不同层：研究、Agent、数据、回测、执行、基础模型")
    groups = [
        ("LLM 因子研发", ["QuantaAlpha", "alpha_mining", "Qlib + RD-Agent"]),
        ("研究与回测框架", ["AKQuant", "Qlib", "vnpy.alpha"]),
        ("行情与信息源", ["a-stock-data"]),
        ("时序基础模型", ["Kronos"]),
        ("交易执行平台", ["vnpy"]),
    ]
    y = 1.6
    colors = [ACCENT, GREEN, YELLOW, RGBColor(186, 104, 200), RGBColor(255, 138, 101)]
    for i, (title, items) in enumerate(groups):
        add_panel(slide, Inches(0.8), Inches(y + i * 1.02), Inches(11.7), Inches(0.82), title)
        add_bullets(
            slide,
            Inches(3.0),
            Inches(y + 0.18 + i * 1.02),
            Inches(9.1),
            Inches(0.38),
            [" / ".join(items)],
            font_size=12,
            color=colors[i],
        )
    add_panel(slide, Inches(0.8), Inches(6.85), Inches(11.7), Inches(0.42), None)
    add_bullets(
        slide,
        Inches(1.0),
        Inches(6.9),
        Inches(11.2),
        Inches(0.25),
        ["结论：USalpha 目前主要覆盖“研究原型 + 页面产品雏形”，但距离“Agent 平台”与“交易平台”仍有明显层级差距。"],
        font_size=11.5,
    )
    add_footer(slide, "06")


def slide_compare_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "参考平台对比（一）", "LLM 因子研发与工业级研究平台")
    add_table(
        slide,
        Inches(0.45),
        Inches(1.35),
        Inches(12.4),
        Inches(5.7),
        ["平台", "核心能力", "USalpha 尚缺内容"],
        [
            ["QuantaAlpha", "LLM 自演化因子挖掘、trajectory 规划、Qlib 回测、实验复现、Web UI", "多轮 Agent 规划、Qlib 级研究底座、系统化实验协议"],
            ["alpha_mining", "Qlib demo、LLM alpha refinement、pairwise ranking、简单前后端", "更正式的平台结构、统一评估体系、排名学习扩展"],
            ["Qlib", "工业级数据层、工作流、模型库、策略/执行/分析、在线服务、RL、PIT 数据", "强数据基础设施、workflow、model zoo、线上化与 nested execution"],
        ],
        col_widths=[Inches(1.6), Inches(4.8), Inches(5.8)],
        font_size=10.5,
    )
    add_footer(slide, "07")


def slide_compare_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "参考平台对比（二）", "多智能体交易、回测执行、数据接入、基础模型")
    add_table(
        slide,
        Inches(0.45),
        Inches(1.35),
        Inches(12.4),
        Inches(5.7),
        ["平台", "核心能力", "USalpha 尚缺内容"],
        [
            ["TradingAgents", "LangGraph 多智能体交易组织、分析师/研究员/交易员/风控/组合经理、checkpoint+memory、多模型供应商", "角色化 Agent 编排、状态机、辩论机制、持久记忆、决策日志"],
            ["AKQuant", "Rust+Python 事件驱动回测、复杂订单、ML/WFO、流式回测、准实盘接口", "统一事件引擎、执行语义、复杂订单、流式监控"],
            ["vnpy", "成熟交易平台、丰富 gateway、算法交易、风险管理、数据记录、vnpy.alpha", "真实 OMS/RMS、网关生态、生产部署能力"],
            ["a-stock-data", "A 股全栈数据 skill：行情、估值、研报、热点、北向、龙虎榜、公告", "数据源广度、研究信息层、估值与题材信号层"],
            ["Kronos", "金融 K 线 foundation model、forecasting、batch predict、finetune pipeline", "时序基础模型、预训练+微调路径、生成式市场表征能力"],
        ],
        col_widths=[Inches(1.55), Inches(4.85), Inches(5.95)],
        font_size=10.0,
    )
    add_footer(slide, "08")


def slide_tradingagents_focus(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "TradingAgents 值得借鉴什么", "它提供的不是 alpha 库，而是多角色研究与决策组织")
    add_panel(slide, Inches(0.6), Inches(1.4), Inches(5.9), Inches(4.9), "TradingAgents 已实现的任务")
    add_bullets(
        slide,
        Inches(0.82),
        Inches(1.78),
        Inches(5.45),
        Inches(4.2),
        [
            "把一次股票判断拆成 analyst -> researcher debate -> trader -> risk debate -> portfolio manager 的链路。",
            "用 LangGraph 维护节点状态、条件跳转、checkpoint resume 与长期 decision memory。",
            "接入新闻、情绪、技术、基本面等多源信息，输出结构化交易决策。",
            "支持多家 LLM 供应商与 CLI 工作流，强调可恢复、可复盘、可替换模型。",
        ],
        font_size=12,
    )
    add_panel(slide, Inches(6.7), Inches(1.4), Inches(5.95), Inches(4.9), "对 USalpha 的启发")
    add_bullets(
        slide,
        Inches(6.92),
        Inches(1.78),
        Inches(5.45),
        Inches(4.2),
        [
            "USalpha 现在更像单研究员工作台，缺少“角色分工 + 争论 + 审核 + 最终批准”结构。",
            "可以借鉴它的 state graph、structured output、checkpoint、memory log，而不是照搬其美股数据流。",
            "在我们的目标里，TradingAgents 最适合作为 agent orchestration 参考，不是 research engine 替代品。",
            "对于 A 股场景，最关键的迁移点是：多角色协作框架 + 可恢复执行 + 决策留痕。",
        ],
        font_size=12,
    )
    add_footer(slide, "09")


def slide_full_trade_chain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "完整交易链应该长什么样", "先建立统一坐标系，再看每个库覆盖了哪一段")
    steps = [
        ("数据", "行情/财务/新闻\nPIT/缓存"),
        ("研究", "因子/特征\n模型/规则"),
        ("信号", "预期收益\n排序/概率"),
        ("组合", "仓位/约束\n调仓计划"),
        ("风控", "限仓/暴露\n审批/拦截"),
        ("执行", "订单/撮合\n算法交易/网关"),
        ("监控", "持仓/延迟\n告警/审计"),
        ("复盘", "绩效归因\n记忆/再训练"),
    ]
    add_flow_steps(
        slide,
        steps,
        left=Inches(0.35),
        top=Inches(1.7),
        width=Inches(12.7),
        height=Inches(1.35),
        gap=0.06,
        title_size=10.0,
        body_size=8.2,
        colors=[PANEL, PANEL, RGBColor(36, 66, 95), RGBColor(36, 66, 95), RGBColor(91, 73, 39), RGBColor(91, 73, 39), RGBColor(64, 54, 88), RGBColor(64, 54, 88)],
    )
    add_table(
        slide,
        Inches(0.55),
        Inches(3.45),
        Inches(12.15),
        Inches(3.0),
        ["环节", "作用原理", "USalpha", "优先参考库"],
        [
            ["数据", "保证研究与交易看到的是可用且可追溯的数据", "部分有", "a-stock-data / Qlib / vnpy"],
            ["研究", "把市场信息转换成可验证的预测结构", "有", "USalpha / QuantaAlpha / alpha_mining / Kronos"],
            ["信号", "把模型输出标准化成可下游消费的分数或标签", "有", "Qlib / vnpy.alpha"],
            ["组合", "把单票判断变成资金分配与换仓决策", "弱", "Qlib / vnpy / AKQuant"],
            ["风控", "在下单前和持仓中控制风险暴露", "无", "AKQuant / vnpy"],
            ["执行", "把目标仓位转成真实订单与成交", "无", "AKQuant / vnpy"],
            ["监控", "保证系统活着、状态一致、异常可见", "很弱", "AKQuant / vnpy / TradingAgents"],
            ["复盘", "把结果沉淀成知识，驱动下一轮优化", "弱", "TradingAgents / Qlib / QuantaAlpha"],
        ],
        col_widths=[Inches(1.1), Inches(4.7), Inches(1.2), Inches(5.15)],
        font_size=9.5,
    )
    add_footer(slide, "05")


def slide_usalpha_coverage(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "USalpha 在完整交易链上的覆盖", "这一页回答：我们到底做到哪一步了")
    add_table(
        slide,
        Inches(0.42),
        Inches(1.32),
        Inches(12.5),
        Inches(5.95),
        ["交易链阶段", "当前状态", "已有实现", "主要缺失", "优先参考库"],
        [
            ["股票池/研究目标", "部分", "页面/CLI 可选 universe", "缺统一 universe registry", "Qlib / vnpy.alpha"],
            ["数据采集与治理", "部分", "US yfinance + A股缓存", "缺 PIT / 质量层 / 公告新闻统一层", "a-stock-data / Qlib"],
            ["因子与特征", "较强", "526 因子 + LLM 候选因子", "缺注册表与依赖图", "QuantaAlpha / alpha_mining"],
            ["模型与信号", "中等", "ridge baseline / 单因子评分", "缺模型族 / WFO / online score", "Qlib / vnpy.alpha / Kronos"],
            ["组合构建", "较弱", "long-short 分桶 / 候选列表", "缺权重约束、容量、换手控制", "Qlib / AKQuant"],
            ["风控与执行计划", "无", "-", "缺 RMS、成本、订单算法", "AKQuant / vnpy"],
            ["OMS/网关/实盘", "无", "-", "缺 paper/live bridge 与 order lifecycle", "vnpy / AKQuant"],
            ["监控/复盘/记忆", "较弱", "artifacts / json / csv", "缺决策日志、告警、状态恢复", "TradingAgents / Qlib / AKQuant"],
        ],
        col_widths=[Inches(1.55), Inches(1.0), Inches(3.1), Inches(3.5), Inches(3.35)],
        font_size=9.4,
    )
    add_footer(slide, "06")


def slide_quantaalpha_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "QuantaAlpha 流程图", "它解决的是“如何系统地产生和筛选 alpha”")
    add_flow_steps(
        slide,
        [
            ("研究方向", "自然语言主题\n如 Price-Volume"),
            ("规划", "diversified planning\n拆多条方向"),
            ("假设生成", "Hypothesis\n提出因子思想"),
            ("代码实现", "LLM coder\n转表达式/代码"),
            ("局部回测", "validation set\n快速筛选"),
            ("演化", "mutation/crossover\n保留好轨迹"),
            ("因子库", "library json\n沉淀结果"),
            ("独立回测", "Qlib full backtest\n样本外验证"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.4,
        body_size=7.9,
        colors=[PANEL, PANEL, PANEL, PANEL, RGBColor(42, 76, 112), RGBColor(42, 76, 112), RGBColor(42, 76, 112), RGBColor(56, 102, 72)],
    )
    add_panel(slide, Inches(0.6), Inches(3.55), Inches(6.0), Inches(2.5), "每一步的原理")
    add_bullets(slide, Inches(0.82), Inches(3.9), Inches(5.55), Inches(1.95), [
        "先把模糊研究意图拆成多条方向，避免单一路径过早收敛。",
        "每轮都要求 hypothesis -> code -> backtest -> feedback 闭环，确保不是纯文本空想。",
        "通过 mutation/crossover 保留有效轨迹，提升探索效率。"
    ], font_size=11.0)
    add_panel(slide, Inches(6.82), Inches(3.55), Inches(5.85), Inches(2.5), "对我们的价值")
    add_bullets(slide, Inches(7.04), Inches(3.9), Inches(5.4), Inches(1.95), [
        "适合补 USalpha 的研究自动化与实验协议层。",
        "不负责真实订单、风控、网关，所以不能当完整交易平台。",
    ], font_size=11.0)
    add_footer(slide, "06")


def slide_alpha_mining_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "alpha_mining 流程图", "它是更轻量的 Qlib 研究闭环 demo")
    add_flow_steps(
        slide,
        [
            ("Qlib 数据", "DatasetH\n训练/验证/测试"),
            ("MyAlpha", "200+ 因子\nhandler"),
            ("训练模型", "LGBModel\n预测未来收益"),
            ("信号记录", "SignalRecord\n保存预测"),
            ("信号分析", "SigAnaRecord\nIC等统计"),
            ("组合回测", "TopkDropout\n调仓回测"),
            ("可视化", "plot / backend\n展示结果"),
            ("LLM迭代", "refine alpha\n继续试验"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.5,
        body_size=7.9,
    )
    add_panel(slide, Inches(0.6), Inches(3.55), Inches(6.0), Inches(2.5), "原理")
    add_bullets(slide, Inches(0.82), Inches(3.9), Inches(5.55), Inches(1.95), [
        "核心是标准 Qlib 路径：dataset -> model -> signal -> strategy -> portfolio analysis。",
        "TopkDropout 的含义是按分数选 topK，并每期替换部分持仓，模拟现实换仓。"
    ], font_size=11.0)
    add_panel(slide, Inches(6.82), Inches(3.55), Inches(5.85), Inches(2.5), "对我们的价值")
    add_bullets(slide, Inches(7.04), Inches(3.9), Inches(5.4), Inches(1.95), [
        "适合借鉴其标准化 signal/record/backtest 结构。",
        "比 QuantaAlpha 简洁，更像 USalpha 可以先走到的下一步。"
    ], font_size=11.0)
    add_footer(slide, "07")


def slide_qlib_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Qlib 流程图", "它覆盖了从研究到在线模拟交易的完整平台骨架")
    add_flow_steps(
        slide,
        [
            ("数据层", "provider / PIT\ncalendar/features"),
            ("信息提取", "DataHandler\nDataset"),
            ("预测模型", "Model zoo\nsignal"),
            ("决策生成", "strategy\nportfolio/orders"),
            ("执行环境", "exchange sim\ncost/slippage"),
            ("分析记录", "Signal/Port\nworkflow records"),
            ("在线生成", "generate order list"),
            ("在线执行更新", "execute -> update\n账户演进"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.3,
        body_size=7.8,
        colors=[PANEL, PANEL, PANEL, RGBColor(36, 66, 95), RGBColor(91, 73, 39), RGBColor(64, 54, 88), RGBColor(64, 54, 88), RGBColor(64, 54, 88)],
    )
    add_panel(slide, Inches(0.6), Inches(3.5), Inches(12.0), Inches(2.6), "原理")
    add_bullets(slide, Inches(0.85), Inches(3.88), Inches(11.45), Inches(2.0), [
        "Qlib 把“预测”与“决策/执行”分开：模型只给分数，策略把分数转为组合和订单。",
        "在线模块按 generate -> execute -> update 三步推进，保证账户状态、交易结果和下一日决策串起来。",
        "支持多层 strategy/executor 嵌套，这是研究平台向执行平台过渡的关键能力。"
    ], font_size=11.0)
    add_footer(slide, "08")


def slide_tradingagents_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "TradingAgents 流程图", "它把交易判断做成多角色组织流程")
    add_flow_steps(
        slide,
        [
            ("Analysts", "technical/news\nsentiment/fundamental"),
            ("工具调用", "market/news\nfundamental tools"),
            ("Bull/Bear", "研究员辩论\n多空观点"),
            ("Research Manager", "汇总成\n投资计划"),
            ("Trader", "形成交易方案\n时机与仓位"),
            ("Risk Debate", "激进/中性/保守\n风控辩论"),
            ("Portfolio Manager", "最终批准\n交易决策"),
            ("Memory/Resume", "decision log\ncheckpoint"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.2,
        body_size=7.8,
    )
    add_panel(slide, Inches(0.6), Inches(3.55), Inches(6.0), Inches(2.5), "原理")
    add_bullets(slide, Inches(0.82), Inches(3.9), Inches(5.55), Inches(1.95), [
        "用 LangGraph 把角色和条件跳转写成状态机，而不是把 prompt 串成一条长链。",
        "辩论的作用不是热闹，而是强制暴露相反证据，减少单模型单视角偏见。"
    ], font_size=11.0)
    add_panel(slide, Inches(6.82), Inches(3.55), Inches(5.85), Inches(2.5), "局限")
    add_bullets(slide, Inches(7.04), Inches(3.9), Inches(5.4), Inches(1.95), [
        "更像研究型决策组织，不是成熟的订单执行与 broker 平台。",
        "最适合作为 USalpha 的 agent 编排层参考。"
    ], font_size=11.0)
    add_footer(slide, "09")


def slide_akquant_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "AKQuant 流程图", "它最接近完整交易系统的工程链路")
    add_flow_steps(
        slide,
        [
            ("DataFeed", "历史/实时\nBar/Tick流"),
            ("EventEngine", "统一事件总线\n时间推进"),
            ("Strategy", "on_bar/on_tick\n生成订单意图"),
            ("RiskManager", "预交易检查\n限仓/T+1/资金"),
            ("ExecutionClient", "撮合/柜台\n成交语义"),
            ("Portfolio", "持仓/现金\n权益更新"),
            ("Warm Start", "checkpoint\n重启恢复"),
            ("Monitor/Live", "日志/告警\n主备/流式"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.1,
        body_size=7.7,
        colors=[PANEL, PANEL, PANEL, RGBColor(91, 73, 39), RGBColor(91, 73, 39), RGBColor(64, 54, 88), RGBColor(64, 54, 88), RGBColor(64, 54, 88)],
    )
    add_panel(slide, Inches(0.6), Inches(3.5), Inches(12.0), Inches(2.6), "原理")
    add_bullets(slide, Inches(0.85), Inches(3.88), Inches(11.45), Inches(2.0), [
        "Strategy 只决定想下什么单，RiskManager 在引擎层独立拦截，ExecutionClient 负责成交语义，Portfolio 负责状态更新。",
        "这种拆法的价值是：研究逻辑、风控逻辑、成交逻辑、账户逻辑互不污染，才适合走向实盘。",
        "Warm Start 和 live orchestration 解决的是‘系统不能死、死了能接着跑’的问题。"
    ], font_size=11.0)
    add_footer(slide, "10")


def slide_vnpy_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "vn.py 流程图", "它是最成熟的实盘平台型参考")
    add_flow_steps(
        slide,
        [
            ("Data/Gateway", "行情接口\n柜台接口"),
            ("Event Engine", "事件驱动\n统一分发"),
            ("App/Strategy", "CTA/组合/脚本\n交易应用"),
            ("OMS", "订单状态机\n报撤单生命周期"),
            ("RMS", "风控规则\n前端限制"),
            ("Algo Trading", "TWAP/Iceberg\n智能拆单"),
            ("Recorder/Manager", "数据落库\n组合统计"),
            ("Web/RPC", "分布式接入\n运维界面"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.1,
        body_size=7.7,
    )
    add_panel(slide, Inches(0.6), Inches(3.5), Inches(12.0), Inches(2.6), "原理")
    add_bullets(slide, Inches(0.85), Inches(3.88), Inches(11.45), Inches(2.0), [
        "核心是‘平台先行’：先有 gateway、event、OMS、RMS、app，再让不同策略往里挂。",
        "vnpy.alpha 则把研究端补上：dataset -> model -> signal -> strategy backtest，与主交易平台衔接。",
        "这类系统的优势是稳定、接口丰富、生产化强，但研发自由度比研究原型低。"
    ], font_size=11.0)
    add_footer(slide, "11")


def slide_a_stock_data_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "a-stock-data 流程图", "它不是交易引擎，而是 A 股研究信息层")
    add_flow_steps(
        slide,
        [
            ("行情层", "K线/盘口\nPE/PB/换手"),
            ("研报层", "研报列表/PDF\n一致预期"),
            ("信号层", "热点/题材\n北向/龙虎榜"),
            ("新闻层", "个股/快讯\n全球资讯"),
            ("基础层", "F10/季报\n公司资料"),
            ("公告层", "沪深北公告"),
            ("研究拼装", "单票调研\n主题追踪"),
            ("给上游消费", "供 Agent/研究\n生成解释"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.1,
        body_size=7.7,
    )
    add_panel(slide, Inches(0.6), Inches(3.55), Inches(12.0), Inches(2.45), "原理")
    add_bullets(slide, Inches(0.85), Inches(3.92), Inches(11.45), Inches(1.85), [
        "它把分散、异构、反爬复杂的 A 股信息端点封装成统一调用层，解决‘研究没数据、Agent 没上下文’的问题。",
        "对 USalpha 的价值不是回测，而是补齐解释性研究材料、题材/资金/公告等弱结构化信息。"
    ], font_size=11.0)
    add_footer(slide, "12")


def slide_kronos_flow(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Kronos 流程图", "它提供的是时序基础模型，不是交易框架")
    add_flow_steps(
        slide,
        [
            ("K线数据", "OHLCV\n多市场"),
            ("Tokenizer", "连续值量化成\n层级离散token"),
            ("预训练", "自回归 Transformer\n学市场语言"),
            ("预测器", "future OHLCV\nbatch forecast"),
            ("微调", "用本地任务\n继续训练"),
            ("Qlib数据预处理", "prepare train/val/test"),
            ("回测示例", "简单策略\n验证效果"),
            ("信号接入", "供研究层消费"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.0,
        body_size=7.6,
    )
    add_panel(slide, Inches(0.6), Inches(3.55), Inches(12.0), Inches(2.45), "原理")
    add_bullets(slide, Inches(0.85), Inches(3.92), Inches(11.45), Inches(1.85), [
        "先把连续 K 线序列 token 化，再用大模型学习时序结构，这样模型学到的是市场形态语言而不是传统手工特征。",
        "对 USalpha 的意义是补一类新模型入口，可作为 model_lab 的高级时序模型选项。"
    ], font_size=11.0)
    add_footer(slide, "13")


def slide_best_of_breed(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "博采众长后的目标流程图", "把每个库最擅长的部分组合起来")
    add_flow_steps(
        slide,
        [
            ("数据底座", "Qlib + a-stock-data\nPIT/行情/题材/公告"),
            ("研究生产", "USalpha + QuantaAlpha\n因子/特征/自动实验"),
            ("高级模型", "Kronos / ML zoo\n时序与截面"),
            ("信号标准化", "Qlib/vnpy.alpha\nscore -> signal"),
            ("组合决策", "Qlib/AKQuant\n仓位与约束"),
            ("Agent 审议", "TradingAgents\n多角色审查"),
            ("风控执行", "AKQuant + vnpy\nRMS/OMS/gateway"),
            ("监控复盘", "TradingAgents memory\n日志/归因/再训练"),
        ],
        left=Inches(0.3),
        top=Inches(1.7),
        width=Inches(12.75),
        height=Inches(1.35),
        gap=0.05,
        title_size=9.1,
        body_size=7.7,
    )
    add_table(
        slide,
        Inches(0.55),
        Inches(3.45),
        Inches(12.15),
        Inches(2.9),
        ["环节", "最优参考", "USalpha 当前状态", "下一步"],
        [
            ["数据", "Qlib + a-stock-data", "弱到中", "先补 catalog / PIT / A股信息层"],
            ["研究自动化", "QuantaAlpha", "中", "补 experiment loop 与 factor library"],
            ["Agent", "TradingAgents", "弱", "引入 graph / memory / reviewer"],
            ["执行", "AKQuant + vnpy", "无", "抽 backtest_core -> OMS/RMS/gateway"],
            ["复盘", "Qlib + TradingAgents", "弱", "统一 artifacts + decision memory"],
        ],
        col_widths=[Inches(1.4), Inches(2.8), Inches(2.0), Inches(5.95)],
        font_size=9.4,
    )
    add_footer(slide, "14")


def slide_signal_to_trade_chain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "信号之后，交易是怎么完成的", "真正难的是把 score 变成仓位、订单、成交和风险约束")
    add_flow_steps(
        slide,
        [
            ("1 信号", "pred / factor / rank\n先得到相对强弱"),
            ("2 选股", "选 topK / 阈值\n确定候选池"),
            ("3 仓位", "等权/市值/优化\n决定买多少"),
            ("4 调仓", "保留/卖出/替换\n决定换多少"),
            ("5 下单", "目标仓位 -> 订单\n数量/价格/时机"),
            ("6 成交", "撮合/滑点/涨跌停\n不一定全成"),
            ("7 风控", "仓位/行业/回撤\n前后都检查"),
            ("8 复盘", "收益/回撤/换手\n归因后再优化"),
        ],
        left=Inches(0.28),
        top=Inches(1.7),
        width=Inches(12.8),
        height=Inches(1.38),
        gap=0.05,
        title_size=9.3,
        body_size=7.8,
        colors=[PANEL, PANEL, RGBColor(36, 66, 95), RGBColor(36, 66, 95), RGBColor(91, 73, 39), RGBColor(91, 73, 39), RGBColor(64, 54, 88), RGBColor(64, 54, 88)],
    )
    add_panel(slide, Inches(0.65), Inches(3.55), Inches(12.0), Inches(2.55), "核心原则")
    add_bullets(slide, Inches(0.9), Inches(3.92), Inches(11.45), Inches(1.95), [
        "信号只回答“谁更好”，交易系统还要回答“买多少、什么时候买、买不买得进、出了风险怎么办”。",
        "平均满仓买入只是最简单的仓位规则之一，真实平台通常会同时控制换手、行业暴露、单票权重、成交量占比和回撤。",
        "所以交易描述至少要拆成四层：选股规则、仓位规则、调仓规则、风控规则。"
    ], font_size=11.0)
    add_footer(slide, "15")


def slide_usalpha_trade_mechanics(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "USalpha 现在到底怎么交易", "当前主流程和 A 股页面回测其实是两套不同层级")
    add_table(
        slide,
        Inches(0.42),
        Inches(1.32),
        Inches(12.5),
        Inches(5.95),
        ["场景", "选股方式", "仓位方式", "调仓方式", "成交假设", "风控现状"],
        [
            ["主研究回测", "按 pred 分数取 top/bottom quantile", "不是实际持仓；直接算多组平均收益", "每个交易日重算分组", "没有真实订单/撮合", "几乎没有，只统计收益与回撤"],
            ["A 股页面回测", "按当日或前一日 score 取 topN", "等权分配现金，100股整数手，接近满仓", "先全部卖出，再平均买入新 topN", "按 open/close 直接成交，默认全成", "无涨跌停过滤、无行业约束、无单票上限、无成交量限制"],
            ["候选股输出", "从 best factor 或模型信号选前几名", "只输出名单，不算权重", "人工看后决定是否换仓", "没有模拟", "完全依赖人工判断"],
        ],
        col_widths=[Inches(1.55), Inches(2.25), Inches(2.05), Inches(2.05), Inches(2.15), Inches(2.45)],
        font_size=9.1,
    )
    add_footer(slide, "16")


def slide_usalpha_trade_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "USalpha 的当前交易细节", "你说的“平均满仓买入”主要对应 A 股页面回测")
    add_panel(slide, Inches(0.6), Inches(1.45), Inches(3.85), Inches(4.9), "当前规则")
    add_bullets(slide, Inches(0.82), Inches(1.8), Inches(3.4), Inches(4.2), [
        "先把旧持仓全部卖掉。",
        "对 topN 股票平均分现金预算。",
        "每只股票按 100 股整数手向下取整。",
        "买不完的钱留在现金里。",
        "下一次调仓再重复一次。",
    ], font_size=12)
    add_panel(slide, Inches(4.72), Inches(1.45), Inches(3.9), Inches(4.9), "这意味着什么")
    add_bullets(slide, Inches(4.94), Inches(1.8), Inches(3.45), Inches(4.2), [
        "本质上是 long-only、等权、接近满仓。",
        "没有保留核心仓位，也没有渐进换仓。",
        "换手率可能很高，交易成本被低估。",
        "如果 topN 都高相关，会集中暴露在同一板块。",
    ], font_size=12)
    add_panel(slide, Inches(8.84), Inches(1.45), Inches(3.9), Inches(4.9), "它缺什么")
    add_bullets(slide, Inches(9.06), Inches(1.8), Inches(3.45), Inches(4.2), [
        "单票权重上限",
        "行业/风格暴露约束",
        "最小持有期与 turnover 控制",
        "涨跌停/停牌/成交量容量约束",
        "止损、回撤线、熔断规则",
    ], font_size=12)
    add_footer(slide, "17")


def slide_qilib_trade_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Qlib 怎么交易", "它把“选股”和“仓位”分成简单版和优化版两套")
    add_panel(slide, Inches(0.6), Inches(1.45), Inches(5.9), Inches(4.9), "TopkDropout")
    add_bullets(slide, Inches(0.82), Inches(1.8), Inches(5.45), Inches(4.2), [
        "先按 score 排名，持有 topK。",
        "每天优先卖掉持仓里排名最差的 n_drop 只，再买入未持仓里最好的替代者。",
        "风险资金比例 `risk_degree` 默认约 95%，所以不是 100% 满仓。",
        "买入资金按新买入股票数均分，天然限制单票权重约为 95%/K。",
        "可设置 `hold_thresh`、`only_tradable`、`forbid_all_trade_at_limit`。"
    ], font_size=11.2)
    add_panel(slide, Inches(6.72), Inches(1.45), Inches(5.95), Inches(4.9), "EnhancedIndexing")
    add_bullets(slide, Inches(6.94), Inches(1.8), Inches(5.45), Inches(4.2), [
        "不再等权，而是做组合优化。",
        "目标是在追踪误差受控前提下，提高主动收益。",
        "会用风险模型、协方差、行业暴露等信息求最优权重。",
        "这类策略更接近机构实盘，而不是简单 topK 轮换。"
    ], font_size=11.2)
    add_footer(slide, "18")


def slide_vnpy_akquant_trade_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "vn.py 和 AKQuant 怎么交易", "它们更强调目标仓位、订单生命周期和风险边界")
    add_table(
        slide,
        Inches(0.42),
        Inches(1.32),
        Inches(12.5),
        Inches(5.95),
        ["平台", "常见仓位规则", "调仓方式", "订单层能力", "风控能力"],
        [
            ["vnpy.alpha", "示例策略是 topK + n_drop + cash_ratio 95% + 最小持有期", "先定 target，再执行 execute_trading", "有买卖、撤单、价格偏移、成交回报、逐日盯市", "手续费、最小佣金、最小持有期，易扩展到更强交易平台层"],
            ["AKQuant", "可用 order_target_percent / order_target_weights / order_target_positions", "按目标权重或目标仓位调仓", "事件驱动、限价、滑点、成交量限制、warm start、live bridge", "max_position_pct、max_order_value、sector_concentration、max_drawdown、stop_loss_threshold、restricted_list"],
        ],
        col_widths=[Inches(1.4), Inches(3.2), Inches(2.2), Inches(2.7), Inches(3.0)],
        font_size=9.2,
    )
    add_footer(slide, "19")


def slide_trading_risk_patterns(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "常见交易策略与风险控制模式", "这页回答：别的库通常怎么做，而不是只做平均满仓")
    add_table(
        slide,
        Inches(0.42),
        Inches(1.32),
        Inches(12.5),
        Inches(5.95),
        ["类别", "典型规则", "优点", "主要风险", "适合参考库"],
        [
            ["等权 topK", "每只股票同样资金，接近满仓", "简单稳健、解释性强", "高换手、行业集中", "USalpha 当前 / Qlib / vnpy"],
            ["TopK+Dropout", "持有 topK，每次只换出 n_drop", "控制换手，不必全换仓", "依然缺显式风险模型", "Qlib / alpha_mining / vnpy"],
            ["目标权重调仓", "直接给每只股票 2%、3% 之类权重", "容易加单票上限和现金比例", "权重设计要合理", "AKQuant / vnpy"],
            ["优化组合", "收益最大化 + 风险暴露约束", "更接近机构组合管理", "对风险模型质量要求高", "Qlib EnhancedIndexing"],
            ["多空中性", "多头和空头同时持有", "削弱市场 beta", "真实执行更复杂", "USalpha 主回测 / Qlib / AKQuant"],
            ["事件驱动执行", "按回报、滑点、成交量逐笔推进", "更接近真实市场", "开发复杂度高", "AKQuant / vnpy"],
        ],
        col_widths=[Inches(1.55), Inches(3.0), Inches(2.05), Inches(2.45), Inches(3.45)],
        font_size=9.1,
    )
    add_footer(slide, "20")


def slide_usalpha_trade_upgrade(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "USalpha 的交易升级建议", "先把平均满仓买入升级成可控的交易规则")
    add_panel(slide, Inches(0.6), Inches(1.45), Inches(3.9), Inches(4.9), "第一步：从全换仓到渐进换仓")
    add_bullets(slide, Inches(0.82), Inches(1.8), Inches(3.45), Inches(4.2), [
        "引入 topK + n_drop。",
        "不是每天全卖全买，而是只替换尾部几只。",
        "保留最强持仓，显著降低 turnover。"
    ], font_size=12)
    add_panel(slide, Inches(4.72), Inches(1.45), Inches(3.9), Inches(4.9), "第二步：从等权到约束权重")
    add_bullets(slide, Inches(4.94), Inches(1.8), Inches(3.45), Inches(4.2), [
        "增加 cash_ratio，例如 90% 或 95%。",
        "设置单票上限，例如 5% 或 10%。",
        "限制行业暴露、风格暴露和最小持有期。"
    ], font_size=12)
    add_panel(slide, Inches(8.84), Inches(1.45), Inches(3.9), Inches(4.9), "第三步：真实风控")
    add_bullets(slide, Inches(9.06), Inches(1.8), Inches(3.45), Inches(4.2), [
        "涨跌停/停牌过滤",
        "成交量参与率限制",
        "最大回撤暂停交易",
        "单日亏损限额",
        "订单与持仓审计"
    ], font_size=12)
    add_footer(slide, "21")


def slide_gaps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "如果目标是 AI Agent 量化平台，USalpha 还缺什么", "缺的不是一个模型，而是一整套平台层")
    add_panel(slide, Inches(0.6), Inches(1.45), Inches(3.9), Inches(4.9), "缺口一：平台基础设施")
    add_bullets(slide, Inches(0.82), Inches(1.8), Inches(3.45), Inches(4.2), [
        "统一数据目录与元数据 catalog",
        "PIT / corporate action / universe 管理",
        "实验 registry / run lineage / 配置管理",
        "统一策略、回测、评估接口",
        "模型注册、因子注册、数据质量检查",
    ], font_size=12)
    add_panel(slide, Inches(4.7), Inches(1.45), Inches(3.9), Inches(4.9), "缺口二：Agent 能力层")
    add_bullets(slide, Inches(4.92), Inches(1.8), Inches(3.45), Inches(4.2), [
        "研究目标拆解与任务规划",
        "因子提出、检索、反思、去重",
        "代码生成与实验执行调度",
        "结果判读、失败归因、版本记忆",
        "报告自动生成与人机协同接口",
    ], font_size=12)
    add_panel(slide, Inches(8.8), Inches(1.45), Inches(3.9), Inches(4.9), "缺口三：交易与运营层")
    add_bullets(slide, Inches(9.02), Inches(1.8), Inches(3.45), Inches(4.2), [
        "统一执行引擎与订单抽象",
        "风控、权限、审计、告警",
        "paper/live bridge",
        "多策略编排与资源调度",
        "监控面板与运维工作流",
    ], font_size=12)
    add_footer(slide, "10")


def slide_target_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "建议的目标架构", "把系统拆成三层：Research / Agent / Execution")
    layers = [
        ("Presentation", "Dashboard / CLI / Report / Notebook / API"),
        ("Agent Orchestration", "Planner / Researcher / Factor Proposer / Evaluator / Reviewer / Memory"),
        ("Research Platform", "Data Catalog / Feature Store / Model Registry / Experiment Tracker / Backtest Engine"),
        ("Execution Platform", "Portfolio / OMS / RMS / Broker Gateway / Monitoring / Alerting"),
        ("Storage", "Market Data / PIT Data / Artifacts / Logs / Vector DB / Knowledge Base"),
    ]
    top = 1.5
    heights = [0.7, 1.0, 1.05, 1.0, 0.8]
    colors = [ACCENT, RGBColor(92, 107, 192), RGBColor(0, 150, 136), RGBColor(255, 152, 0), PANEL_2]
    for i, (name, desc) in enumerate(layers):
        y = top + sum(heights[:i]) + i * 0.12
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(y), Inches(11.55), Inches(heights[i]))
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i]
        shape.line.fill.background()
        tb = slide.shapes.add_textbox(Inches(1.15), Inches(y + 0.08), Inches(11.0), Inches(heights[i] - 0.12))
        tf = tb.text_frame
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = name
        r1.font.size = Pt(15)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = desc
        r2.font.size = Pt(11)
        r2.font.color.rgb = WHITE
    add_footer(slide, "11")


def slide_modules(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "建议的模块化切分", "按完整交易链拆，而不是按页面功能拆")
    add_table(
        slide,
        Inches(0.45),
        Inches(1.35),
        Inches(12.4),
        Inches(5.8),
        ["模块", "职责", "与当前 USalpha 的关系"],
        [
            ["data_core", "多数据源、缓存、PIT、公告/新闻/题材、元数据 catalog", "从 data.py 与 CN/A股缓存逻辑外扩"],
            ["research_core", "因子表达、特征工程、候选库、研究配置、实验协议", "整合 factors.py / factor_evolution.py"],
            ["model_lab", "训练、预测、滚动训练、模型注册、信号标准化", "扩展当前 model.py"],
            ["portfolio_core", "仓位、约束、换手、容量、组合回测", "从当前 long-short 回测上移一层"],
            ["execution_core", "RMS、OMS、撮合、paper/live gateway", "当前基本缺失，需要新建"],
            ["agent_lab", "planner / reviewer / PM / memory / checkpoint", "把单轮 LLM 工具升级为图编排层"],
            ["ops_runtime", "artifact、日志、监控、告警、恢复、审计", "把零散输出提升为运维层"],
            ["ui_app", "dashboard、CLI、报告、API", "把 factor_dashboard.py 拆成多个入口"],
        ],
        col_widths=[Inches(1.65), Inches(4.5), Inches(6.1)],
        font_size=10.0,
    )
    add_footer(slide, "15")


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "实施路线图", "先补交易链骨架，再上 multi-agent")
    add_panel(slide, Inches(0.6), Inches(1.55), Inches(3.9), Inches(4.7), "阶段 1：研究链硬化")
    add_bullets(slide, Inches(0.82), Inches(1.9), Inches(3.45), Inches(4.0), [
        "统一 data_core 与 A 股信息层",
        "统一 signal / portfolio / backtest 口径",
        "引入 experiment registry / run metadata",
        "把 dashboard 与 service 解耦",
    ], font_size=12)
    add_panel(slide, Inches(4.72), Inches(1.55), Inches(3.9), Inches(4.7), "阶段 2：Agent 化决策")
    add_bullets(slide, Inches(4.94), Inches(1.9), Inches(3.45), Inches(4.0), [
        "planner / analyst / reviewer / PM graph",
        "研究计划、失败归因、知识库记忆",
        "自动报告与人工审批节点",
        "多模型、多数据源协同",
    ], font_size=12)
    add_panel(slide, Inches(8.84), Inches(1.55), Inches(3.9), Inches(4.7), "阶段 3：执行平台化")
    add_bullets(slide, Inches(9.06), Inches(1.9), Inches(3.45), Inches(4.0), [
        "paper/live bridge",
        "OMS / RMS / alerting / audit",
        "多策略编排与资源调度",
        "生产部署与监控体系",
    ], font_size=12)
    add_footer(slide, "18")


def slide_multi_agent_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "USalpha 如何增加 Multi-Agent", "Agent 不是页面助手，而是交易链中的协作层")
    add_panel(slide, Inches(0.6), Inches(1.45), Inches(4.0), Inches(4.9), "放在哪里")
    add_bullets(slide, Inches(0.82), Inches(1.82), Inches(3.52), Inches(4.1), [
        "新增 `agent_lab` 作为独立层。",
        "由 graph/state machine 驱动，而不是由 Streamlit 页面直接串函数。",
        "Agent 只调用标准化 service：data_core / research_core / model_lab / portfolio_core / execution_core。",
        "每次 run 产出 plan、signal review、portfolio verdict、memory。",
    ], font_size=12)
    add_panel(slide, Inches(4.72), Inches(1.45), Inches(4.0), Inches(4.9), "它能实现什么")
    add_bullets(slide, Inches(4.94), Inches(1.82), Inches(3.52), Inches(4.1), [
        "自动提出候选因子与策略假设。",
        "自动拉取数据、运行实验、评估组合、形成交易建议。",
        "对失败结果做归因：数据问题、过拟合、交易成本、样本漂移。",
        "自动输出日报、周报、版本对比和风控提示。",
    ], font_size=12)
    add_panel(slide, Inches(8.84), Inches(1.45), Inches(3.9), Inches(4.9), "工程要求")
    add_bullets(slide, Inches(9.06), Inches(1.82), Inches(3.42), Inches(4.1), [
        "结构化输入输出",
        "checkpoint / resume",
        "memory / knowledge base",
        "任务队列与并发控制",
        "人工批准节点",
    ], font_size=12)
    add_footer(slide, "16")


def slide_multi_agent_roles(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "建议的 Agent 角色设计", "覆盖研究、组合、风控和执行审议")
    add_table(
        slide,
        Inches(0.45),
        Inches(1.35),
        Inches(12.4),
        Inches(5.75),
        ["角色", "主要职责", "调用模块 / 产物"],
        [
            ["Planner", "拆解研究目标、安排实验序列、定义成功标准", "agent graph / run manifest"],
            ["Data Steward", "选择股票池、校验数据质量、补缓存、控制样本口径", "data_core / universe report"],
            ["Factor Proposer", "提出新因子、新组合逻辑、新筛选条件", "research_core / candidate specs"],
            ["Research Coder", "把假设转成可执行配置或代码模板", "research_core / model_lab / patch or config"],
            ["Backtest Runner", "批量执行实验、采集指标、保留 artifact", "portfolio_core / metrics bundle"],
            ["Reviewer", "检查稳健性、过拟合、泄露、成本敏感性", "evaluation report / veto"],
            ["Portfolio Agent", "把单票信号转成仓位、换仓、容量约束", "portfolio plan / target weights"],
            ["Risk Agent", "评估回撤、行业暴露、风格偏移和下单风险", "risk report / limits"],
            ["Execution Agent", "决定 paper/live、下单批次、是否需要算法执行", "order intent / execution plan"],
            ["PM Agent", "综合建议，决定保留、淘汰、上线观察", "final verdict / next action"],
        ],
        col_widths=[Inches(1.55), Inches(4.85), Inches(5.95)],
        font_size=9.6,
    )
    add_footer(slide, "17")


def slide_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "结论与建议", "下一阶段的重点，不是更会找信号，而是更会把信号交易出去")
    add_panel(slide, Inches(0.7), Inches(1.6), Inches(12.0), Inches(4.9), "最终判断")
    add_bullets(slide, Inches(0.95), Inches(1.95), Inches(11.5), Inches(4.2), [
        "USalpha 当前最真实的交易规则是：A 股页面回测里，按 topN 等权、接近满仓、100股整数手、全卖全买；这能跑，但太粗。",
        "Qlib 给出的是 `topK + n_drop + risk_degree` 的简洁组合规则，vnpy 和 AKQuant 则进一步把目标仓位、订单状态、撮合、风控和回撤控制补全。",
        "因此下一阶段最值得做的不是继续提升单个因子，而是把交易规则升级成：渐进换仓、现金比例、单票上限、行业约束、成交量限制、回撤熔断。",
        "只有交易层补上，USalpha 才会从‘研究工作台’走向‘可执行的量化平台’。 ",
    ], font_size=13)
    add_footer(slide, "22")


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs)
    slide_exec_summary(prs)
    slide_usalpha_done(prs)
    slide_full_trade_chain(prs)
    slide_usalpha_flow(prs)
    slide_usalpha_coverage(prs)
    slide_quantaalpha_flow(prs)
    slide_alpha_mining_flow(prs)
    slide_qlib_flow(prs)
    slide_tradingagents_flow(prs)
    slide_akquant_flow(prs)
    slide_vnpy_flow(prs)
    slide_a_stock_data_flow(prs)
    slide_kronos_flow(prs)
    slide_best_of_breed(prs)
    slide_signal_to_trade_chain(prs)
    slide_usalpha_trade_mechanics(prs)
    slide_usalpha_trade_detail(prs)
    slide_qilib_trade_detail(prs)
    slide_vnpy_akquant_trade_detail(prs)
    slide_trading_risk_patterns(prs)
    slide_usalpha_trade_upgrade(prs)
    slide_modularity(prs)
    slide_modules(prs)
    slide_multi_agent_arch(prs)
    slide_multi_agent_roles(prs)
    slide_roadmap(prs)
    slide_closing(prs)
    for i, slide in enumerate(prs.slides, start=1):
        overwrite_footer(slide, f"{i:02d}")
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
